"""Step 4: PPT Builder — fill template with generated content."""

from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Pt, Emu
from pptx.dml.color import RGBColor


def _set_text_with_black_font(
    shape, text: str, font_size_pt: int = 14, alignment: PP_ALIGN | None = None
):
    """Set a shape's first paragraph text and force font color to black.

    Uses a smaller font and disables word wrap so names/affiliations/date
    stay on a single line. The text box width is expanded to fit the content.
    """
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = text
    if alignment is not None:
        p.alignment = alignment
    for run in p.runs:
        run.font.color.rgb = RGBColor(0, 0, 0)
        run.font.size = Pt(font_size_pt)
        run.font.name = "微软雅黑"

    # Heuristic width per character at 14pt: CJK ~0.52 cm, Latin/digit ~0.28 cm
    # Add 1.5 cm padding so the text doesn't feel cramped and fits inside margins
    estimated_width = sum(0.52 if "\u4e00" <= ch <= "\u9fff" else 0.28 for ch in text)
    target_width = min(max(estimated_width + 1.5, 4.0), 18.0)
    shape.width = Cm(target_width)


def _set_title_text(shape, text: str):
    """Set title text: black, centered, with word wrap enabled."""
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    for run in p.runs:
        run.font.color.rgb = RGBColor(0, 0, 0)
        run.font.name = "微软雅黑"

from .base import BaseSkill
from ..models import PresentationContent, SlideContent
from ..utils.logger import logger


class PPTBuilderSkill(BaseSkill):
    """Fill PPT template with generated slide content.

    Strategy:
        1. Load template.pptx
        2. Fill fixed areas (title, subtitle, presenter, date) by shape name matching
        3. For content slides (4-15), use two-column layout:
           - Left: centered text content
           - Right: extracted image (if specified)
        4. Delete unused slides based on total_slides
        5. Save output.pptx
    """

    def __init__(self, job_id: str, work_dir: Path, config: dict):
        super().__init__(job_id, work_dir, config)
        self.settings = config.get("settings", None)

    @property
    def skill_name(self) -> str:
        return "ppt_builder"

    @property
    def output_path(self) -> Path:
        return self.work_dir / "06_output.pptx"

    def execute(self, presentation: PresentationContent, **inputs) -> Path:
        template_path = Path(self.settings.template_path) if self.settings else Path("pptx_template/模板.pptx")
        logger.info(f"Building PPT from template: {template_path}")

        prs = Presentation(str(template_path))

        # Build a mapping from slide_index (1-based) to SlideContent
        slide_map = {s.slide_index: s for s in presentation.slides}

        # Process each slide in the template
        for i, slide in enumerate(prs.slides, 1):
            if i not in slide_map:
                continue
            content = slide_map[i]
            self._fill_slide(slide, content, presentation)

        # Remove unused slides
        self._remove_unused_slides(prs, set(slide_map.keys()))

        # Save
        self.output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(self.output_path))
        logger.info(f"PPT saved: {self.output_path}")
        return self.output_path

    def _fill_slide(self, slide, content: SlideContent, presentation: PresentationContent):
        """Fill a single slide with content."""
        # Find all shapes by name for quick lookup
        shapes_by_name: dict[str, list] = {}
        for shape in slide.shapes:
            shapes_by_name.setdefault(shape.name, []).append(shape)

        # Slide 1: Cover page
        # Layout: left=presenter name, right=affiliation (same row), bottom-center=date
        if content.slide_type == "title":
            if presentation.english_title:
                for shape in shapes_by_name.get("矩形 3", []):
                    _set_title_text(shape, presentation.english_title)
            if content.title:
                for shape in shapes_by_name.get("矩形 1", []):
                    _set_title_text(shape, content.title)

            textboxes_12 = shapes_by_name.get("TextBox 12", [])
            if textboxes_12:
                # Sort by top coordinate so we can distinguish the same row vs bottom row
                textboxes_12_sorted = sorted(textboxes_12, key=lambda s: s.top)
                # The two boxes with nearly the same top are presenter (left) and affiliation (right)
                # The box with a larger top is the date box
                same_row = []
                bottom_row = []
                if len(textboxes_12_sorted) >= 2:
                    # tolerance for top difference ~0.3 cm in EMUs
                    threshold = 200000
                    same_row = [textboxes_12_sorted[0]]
                    for s in textboxes_12_sorted[1:]:
                        if abs(s.top - same_row[0].top) <= threshold:
                            same_row.append(s)
                        else:
                            bottom_row.append(s)
                    same_row.sort(key=lambda s: s.left)
                    bottom_row.sort(key=lambda s: s.left)

                # Fill same row: left = presenter, right = affiliation
                if len(same_row) >= 1:
                    name = getattr(self.settings, "presenter_name", "")
                    _set_text_with_black_font(same_row[0], f"汇报人：{name}" if name else "")
                if len(same_row) >= 2:
                    aff = getattr(self.settings, "presenter_affiliation", "")
                    _set_text_with_black_font(same_row[1], aff)
                # Fill bottom row: date
                for box in bottom_row:
                    date = getattr(self.settings, "presenter_date", "")
                    _set_text_with_black_font(box, f"日期：{date}" if date else "")

        # Slide 2: TOC
        elif content.slide_type == "toc":
            for shape in shapes_by_name.get("内容占位符 2", []):
                if content.bullet_points:
                    tf = shape.text_frame
                    tf.clear()
                    for i, point in enumerate(content.bullet_points):
                        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                        p.text = point
                        p.level = 0

        # Slide 3: Paper info
        elif content.slide_type == "paper_info":
            for shape in shapes_by_name.get("TextBox 8", []):
                shape.text_frame.paragraphs[0].text = content.title
            for shape in shapes_by_name.get("TextBox 12", []):
                if content.bullet_points:
                    shape.text_frame.paragraphs[0].text = "\n".join(content.bullet_points)
            # If a title-page image is provided, add it on the right side
            if content.image_path:
                self._add_right_image(slide, content)

        # Content slides (4-15): two-column layout
        elif content.slide_type in ("background", "objective", "results", "methods", "discussion"):
            # Update title
            for shape in shapes_by_name.get("TextBox 8", []):
                shape.text_frame.paragraphs[0].text = content.title

            # Add left text box (centered) and right image
            if content.bullet_points or content.image_path:
                self._add_content_two_column(slide, content)

    def _add_content_two_column(self, slide, content: SlideContent):
        """Add a two-column layout: left text (centered), right image."""
        # Find title position to determine content top
        title_shape = None
        for shape in slide.shapes:
            if shape.name == "TextBox 8":
                title_shape = shape
                break

        # Content area starts below title
        if title_shape:
            content_top = title_shape.top + title_shape.height + Emu(180000)  # 0.5cm gap
        else:
            content_top = Cm(2.0)

        # Slide dimensions (standard 16:9)
        slide_width = Cm(25.0)
        slide_height = Cm(14.06)

        # Left column: text content (centered)
        left_col_left = Cm(1.5)
        left_col_width = Cm(13.0)
        content_height = slide_height - content_top - Cm(1.0)  # 1cm bottom margin

        if content.bullet_points:
            txBox = slide.shapes.add_textbox(left_col_left, content_top, left_col_width, content_height)
            tf = txBox.text_frame
            tf.word_wrap = True

            for i, point in enumerate(content.bullet_points):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = point
                p.level = 0
                p.alignment = PP_ALIGN.CENTER  # Center alignment
                # Set font
                for run in p.runs:
                    run.font.size = Pt(18)
                    run.font.name = "微软雅黑"

        # Right column: image
        if content.image_path:
            img_path = self.work_dir / "images" / content.image_path
            if not img_path.exists():
                # Fallback: try without images subdir
                img_path = self.work_dir / content.image_path

            if img_path.exists():
                right_col_left = Cm(15.5)
                right_col_width = Cm(8.5)
                right_col_height = content_height

                try:
                    from PIL import Image as PILImage
                    with PILImage.open(img_path) as img:
                        img_w, img_h = img.size
                        aspect = img_w / img_h

                    # Fit image within right column while preserving aspect ratio
                    max_w = right_col_width
                    max_h = right_col_height
                    if aspect > max_w / max_h:
                        pic_width = max_w
                        pic_height = max_w / aspect
                    else:
                        pic_height = max_h
                        pic_width = max_h * aspect

                    # Center the image in the right column
                    pic_left = right_col_left + (right_col_width - pic_width) / 2
                    pic_top = content_top + (right_col_height - pic_height) / 2

                    slide.shapes.add_picture(
                        str(img_path),
                        pic_left,
                        pic_top,
                        width=pic_width,
                        height=pic_height,
                    )
                    logger.debug(f"Added image {content.image_path} to slide")

                    # Add image caption below if available
                    if content.image_caption:
                        caption_top = pic_top + pic_height + Emu(72000)  # 0.2cm gap
                        caption_box = slide.shapes.add_textbox(
                            right_col_left, caption_top, right_col_width, Cm(1.0)
                        )
                        caption_tf = caption_box.text_frame
                        caption_tf.word_wrap = True
                        p = caption_tf.paragraphs[0]
                        p.text = content.image_caption
                        p.alignment = PP_ALIGN.CENTER
                        for run in p.runs:
                            run.font.size = Pt(12)
                            run.font.name = "微软雅黑"
                            run.font.italic = True

                except Exception as e:
                    logger.warning(f"Failed to add image {content.image_path}: {e}")
            else:
                logger.warning(f"Image not found: {img_path}")

    def _add_right_image(self, slide, content: SlideContent):
        """Add an image fitted to the right half of the slide.

        Used for paper_info and other slides that already have a left textbox.
        """
        img_path = self.work_dir / "images" / content.image_path
        if not img_path.exists():
            img_path = self.work_dir / content.image_path
        if not img_path.exists():
            logger.warning(f"Image not found: {img_path}")
            return

        # Use fixed right-column area, centered vertically below title
        right_col_left = Cm(14.0)
        right_col_width = Cm(10.0)
        content_top = Cm(2.5)
        content_height = Cm(10.5)

        try:
            from PIL import Image as PILImage

            with PILImage.open(img_path) as img:
                img_w, img_h = img.size
                aspect = img_w / img_h

            max_w = right_col_width
            max_h = content_height
            if aspect > max_w / max_h:
                pic_width = max_w
                pic_height = max_w / aspect
            else:
                pic_height = max_h
                pic_width = max_h * aspect

            pic_left = right_col_left + (right_col_width - pic_width) / 2
            pic_top = content_top + (content_height - pic_height) / 2

            slide.shapes.add_picture(
                str(img_path),
                pic_left,
                pic_top,
                width=pic_width,
                height=pic_height,
            )
            logger.debug(f"Added right-side image {content.image_path} to slide")

            if content.image_caption:
                caption_top = pic_top + pic_height + Emu(72000)
                caption_box = slide.shapes.add_textbox(
                    right_col_left, caption_top, right_col_width, Cm(1.0)
                )
                caption_tf = caption_box.text_frame
                caption_tf.word_wrap = True
                p = caption_tf.paragraphs[0]
                p.text = content.image_caption
                p.alignment = PP_ALIGN.CENTER
                for run in p.runs:
                    run.font.size = Pt(12)
                    run.font.name = "微软雅黑"
                    run.font.italic = True

        except Exception as e:
            logger.warning(f"Failed to add right image {content.image_path}: {e}")

    def _remove_unused_slides(self, prs, keep_indices: set[int]):
        """Remove slides whose 1-based index is not in keep_indices.

        python-pptx does not expose slide deletion, so we manipulate the
        underlying ``<p:sldIdLst>`` element directly. Deletion must proceed
        from high index to low index so remaining indices stay valid.
        """
        total = len(prs.slides)
        to_remove = [i for i in range(1, total + 1) if i not in keep_indices]

        for idx in reversed(to_remove):
            # idx is 1-based; convert to 0-based for list access
            rId = prs.slides._sldIdLst[idx - 1].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[idx - 1]
            logger.debug(f"Removed slide {idx}")

        return prs
